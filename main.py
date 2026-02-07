from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from pptx import Presentation
from io import BytesIO
import json

app = FastAPI()

@app.post("/rebuild-ppt")
async def rebuild_ppt(file: UploadFile = File(...), content: str = Form(...)):
    new_data = json.loads(content)
    # Open the original file as a template
    prs = Presentation(BytesIO(await file.read()))

    for i, slide in enumerate(prs.slides):
        if i < len(new_data):
            slide_update = new_data[i]
            for shape in slide.shapes:
                if not shape.has_text_frame: 
                    continue
                
                tf = shape.text_frame
                # Ensure we have at least one paragraph to work with
                if not tf.paragraphs:
                    continue
                
                p = tf.paragraphs[0]
                
                # Determine the text to inject based on the shape type
                if shape == slide.shapes.title:
                    new_text = slide_update.get('new_title', '')
                else:
                    new_text = "\n".join(slide_update.get('new_content', []))

                # SURGICAL UPDATE: Use the first run to keep formatting
                if p.runs:
                    p.runs[0].text = new_text
                    # IMPORTANT: Clear all other runs to prevent duplicate/old text
                    for r_idx in range(1, len(p.runs)):
                        p.runs[r_idx].text = ""
                else:
                    # Fallback: if the shape was empty, add a run (will use default style)
                    p.add_run().text = new_text

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return StreamingResponse(
        output, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=Rebuilt_Presentation.pptx"}
    )